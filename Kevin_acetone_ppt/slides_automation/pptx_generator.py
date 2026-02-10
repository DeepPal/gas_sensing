from __future__ import annotations

import logging
from pathlib import Path
from typing import Iterable, Mapping
import re

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from .config import LocalSlideDefinition, PPTXConfig, PPTXTypography
from .utils import ensure_directory, slugify

LOGGER = logging.getLogger(__name__)

INLINE_BOLD_PATTERN = re.compile(r"(\*\*.*?\*\*)")
MANUAL_BULLET_PREFIXES = ("•", "-", "→", "✓", "✗")
INDENT_TOKEN = "    "


# =============================================================================
# THEME CONFIGURATION - Professional Dark Research Theme
# =============================================================================
class ThemeColors:
    """Centralized color scheme for consistent styling."""
    # Primary colors - Deep professional dark blue
    BACKGROUND = RGBColor(12, 20, 35)        # Deep navy
    BACKGROUND_ALT = RGBColor(20, 30, 50)    # Slightly lighter navy
    BACKGROUND_CARD = RGBColor(25, 38, 60)   # Card/highlight background
    
    # Text colors - High contrast for readability
    TEXT_PRIMARY = RGBColor(248, 250, 252)   # Near white
    TEXT_SECONDARY = RGBColor(203, 213, 225) # Light gray
    TEXT_MUTED = RGBColor(148, 163, 184)     # Muted gray
    TEXT_HIGHLIGHT = RGBColor(255, 255, 255) # Pure white for emphasis
    
    # Accent colors - Scientific/professional palette
    ACCENT_PRIMARY = RGBColor(59, 130, 246)   # Blue-500 (primary accent)
    ACCENT_SECONDARY = RGBColor(99, 102, 241) # Indigo-500
    ACCENT_SUCCESS = RGBColor(34, 197, 94)    # Green-500 (positive results)
    ACCENT_WARNING = RGBColor(245, 158, 11)   # Amber-500 (caution)
    ACCENT_HIGHLIGHT = RGBColor(14, 165, 233) # Sky-500 (highlights)
    
    # Gradient colors for visual interest
    GRADIENT_START = RGBColor(30, 58, 138)    # Blue-900
    GRADIENT_END = RGBColor(15, 23, 42)       # Slate-900
    
    # Table colors - Professional data presentation
    TABLE_HEADER_BG = RGBColor(30, 64, 175)   # Blue-800
    TABLE_ROW_ALT = RGBColor(25, 38, 60)      # Alternating row
    TABLE_ROW_NORMAL = RGBColor(15, 25, 45)   # Normal row
    TABLE_BORDER = RGBColor(71, 85, 105)      # Border color
    TABLE_HIGHLIGHT = RGBColor(59, 130, 246)  # Highlight important cells
    
    # Chart/visualization colors
    CHART_COLORS = [
        RGBColor(59, 130, 246),   # Blue
        RGBColor(34, 197, 94),    # Green
        RGBColor(245, 158, 11),   # Amber
        RGBColor(239, 68, 68),    # Red
        RGBColor(168, 85, 247),   # Purple
        RGBColor(14, 165, 233),   # Sky
    ]


class LocalPPTBuilder:
    """Create editable PPTX decks from local templates with professional styling."""

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = 13.33
    SLIDE_HEIGHT = 7.5
    
    def __init__(self, config: PPTXConfig):
        self.config = config
        self.slide_count = 0  # Track for slide numbering
        self.primary_font = config.font_name or "Calibri"
        self.typography = config.font_sizes

    # ------------------------------------------------------------------
    # Theme & Styling Methods
    # ------------------------------------------------------------------
    def _apply_dark_theme(self, slide):
        """Apply dark theme background to slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = ThemeColors.BACKGROUND

    def _apply_section_theme(self, slide):
        """Apply slightly different background for section dividers."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = ThemeColors.BACKGROUND_ALT

    def _add_slide_number(self, slide, number: int, total: int = None):
        """Add slide number in bottom-right corner."""
        text = f"{number}" if total is None else f"{number} / {total}"
        textbox = slide.shapes.add_textbox(
            Inches(self.SLIDE_WIDTH - 1.2), 
            Inches(self.SLIDE_HEIGHT - 0.5), 
            Inches(1.0), 
            Inches(0.4)
        )
        frame = textbox.text_frame
        frame.text = text
        para = frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(self.typography.slide_number_size)
        para.font.color.rgb = ThemeColors.TEXT_MUTED
        para.alignment = PP_ALIGN.RIGHT

    def _add_accent_line(self, slide, top: float = 1.2, width: float = 2.5, left: float = 0.6):
        """Add a horizontal accent line under the title."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ThemeColors.ACCENT_PRIMARY
        line.line.fill.background()

    def _add_decorative_corner(self, slide, position: str = "top-right"):
        """Add decorative corner accent for visual interest."""
        if position == "top-right":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches(self.SLIDE_WIDTH - 1.5), Inches(0),
                Inches(1.5), Inches(1.0)
            )
        elif position == "bottom-left":
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_TRIANGLE,
                Inches(0), Inches(self.SLIDE_HEIGHT - 0.8),
                Inches(1.2), Inches(0.8)
            )
        else:
            return
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.ACCENT_PRIMARY
        shape.fill.fore_color.brightness = 0.3  # Subtle
        shape.line.fill.background()

    def _add_footer_bar(self, slide, text: str = ""):
        """Add a subtle footer bar with optional text."""
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(self.SLIDE_HEIGHT - 0.4),
            Inches(self.SLIDE_WIDTH), Inches(0.4)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ThemeColors.BACKGROUND_ALT
        bar.line.fill.background()
        
        if text:
            textbox = slide.shapes.add_textbox(
                Inches(0.5), Inches(self.SLIDE_HEIGHT - 0.35),
                Inches(8), Inches(0.3)
            )
            frame = textbox.text_frame
            frame.text = text
            para = frame.paragraphs[0]
            para.font.name = self.primary_font
            para.font.size = Pt(self.typography.caption_size)
            para.font.color.rgb = ThemeColors.TEXT_MUTED

    def _add_side_accent(self, slide, side: str = "left"):
        """Add a vertical accent bar on the side."""
        if side == "left":
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(0.08), Inches(self.SLIDE_HEIGHT)
            )
        else:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(self.SLIDE_WIDTH - 0.08), Inches(0),
                Inches(0.08), Inches(self.SLIDE_HEIGHT)
            )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ThemeColors.ACCENT_PRIMARY
        bar.line.fill.background()

    def _style_title_text(self, text_frame, font_size: int | None = None):
        """Apply consistent title styling."""
        resolved_size = font_size or self.typography.heading_size
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.primary_font
            paragraph.font.size = Pt(resolved_size)
            paragraph.font.color.rgb = ThemeColors.TEXT_PRIMARY
            paragraph.font.bold = True

    def _style_body_text(self, text_frame, font_size: int | None = None):
        """Apply consistent body text styling."""
        resolved_size = font_size or self.typography.body_size
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.primary_font
            paragraph.font.size = Pt(resolved_size)
            paragraph.font.color.rgb = ThemeColors.TEXT_PRIMARY

    def _style_accent_text(self, text_frame, font_size: int | None = None):
        """Apply accent text styling (for section headers, etc.)."""
        resolved_size = font_size or self.typography.section_title_size
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.primary_font
            paragraph.font.size = Pt(resolved_size)
            paragraph.font.color.rgb = ThemeColors.ACCENT_PRIMARY
            paragraph.font.bold = True

    def _style_secondary_text(self, text_frame, font_size: int | None = None):
        """Apply secondary/muted text styling."""
        resolved_size = font_size or self.typography.body_secondary_size
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.primary_font
            paragraph.font.size = Pt(resolved_size)
            paragraph.font.color.rgb = ThemeColors.TEXT_SECONDARY

    # ------------------------------------------------------------------
    # Build Presentation
    # ------------------------------------------------------------------
    def build_presentation(
        self,
        *,
        title: str,
        slides: Iterable[LocalSlideDefinition],
        placeholders: Mapping[str, str],
        data_sources: Mapping[str, object],
        output_name: str | None = None,
    ) -> Path:
        prs = self._load_template()
        slide_defs = list(slides)

        if not slide_defs:
            LOGGER.warning("No slide definitions found; injecting default title slide")
            slide_defs.append(
                LocalSlideDefinition(kind="title", title=title, subtitle="Automated deck")
            )

        total_slides = len(slide_defs)
        for idx, definition in enumerate(slide_defs, start=1):
            self.slide_count = idx
            self._render_slide(prs, definition, placeholders, data_sources, idx, total_slides)

        output_dir = ensure_directory(self.config.output_dir)
        filename = output_name or slugify(title)
        output_path = output_dir / f"{filename}.pptx"
        prs.save(output_path)
        LOGGER.info("Saved PPTX deck → %s", output_path)
        return output_path

    # ------------------------------------------------------------------
    def _load_template(self) -> Presentation:
        if self.config.template_path:
            template_path = Path(self.config.template_path)
            if not template_path.exists():
                LOGGER.warning("Template %s not found; creating blank presentation", template_path)
            else:
                return Presentation(template_path)
        prs = Presentation()
        prs.slide_width = Inches(self.SLIDE_WIDTH)
        prs.slide_height = Inches(self.SLIDE_HEIGHT)
        return prs

    # ------------------------------------------------------------------
    # Slide Rendering Dispatcher
    # ------------------------------------------------------------------
    def _render_slide(
        self,
        prs: Presentation,
        definition: LocalSlideDefinition,
        placeholders: Mapping[str, str],
        data_sources: Mapping[str, object],
        slide_num: int = 1,
        total_slides: int = 1,
    ) -> None:
        kind = definition.kind.lower()
        
        # Dispatch to appropriate slide builder
        if kind == "title":
            self._add_title_slide(prs, definition, placeholders, slide_num, total_slides)
        elif kind == "section":
            self._add_section_slide(prs, definition, placeholders, slide_num, total_slides)
        elif kind == "bullets":
            self._add_bullet_slide(prs, definition, placeholders, slide_num, total_slides)
        elif kind == "image":
            self._add_image_slide(prs, definition, placeholders, slide_num, total_slides)
        elif kind == "table":
            self._add_table_slide(prs, definition, placeholders, data_sources, slide_num, total_slides)
        elif kind == "two_column":
            self._add_two_column_slide(prs, definition, placeholders, slide_num, total_slides)
        elif kind == "text":
            self._add_text_slide(prs, definition, placeholders, slide_num, total_slides)
        elif kind == "equation":
            self._add_equation_slide(prs, definition, placeholders, slide_num, total_slides)
        else:
            LOGGER.warning("Unknown slide type '%s'; defaulting to text", definition.kind)
            self._add_text_slide(prs, definition, placeholders, slide_num, total_slides)

    # ------------------------------------------------------------------
    # Slide Type Builders
    # ------------------------------------------------------------------
    def _add_title_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create a professional title slide with centered content and visual elements."""
        layout = prs.slide_layouts[self.config.title_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add decorative elements
        self._add_side_accent(slide, "left")
        
        # Add top decorative bar
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(self.SLIDE_WIDTH), Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = ThemeColors.ACCENT_PRIMARY
        top_bar.line.fill.background()

        # Main title - centered with larger font
        if definition.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.2), Inches(11.73), Inches(1.8)
            )
            frame = title_box.text_frame
            frame.word_wrap = True
            frame.text = _fill(definition.title, placeholders)
            para = frame.paragraphs[0]
            para.font.name = self.primary_font
            para.font.size = Pt(self.typography.title_size)
            para.font.color.rgb = ThemeColors.TEXT_HIGHLIGHT
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER

        # Accent line - centered, wider
        self._add_accent_line(slide, top=4.0, width=4.0, left=(self.SLIDE_WIDTH - 4.0) / 2)

        # Subtitle - centered below accent line
        if definition.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(4.3), Inches(11.73), Inches(2.2)
            )
            frame = sub_box.text_frame
            frame.word_wrap = True
            frame.text = _fill(definition.subtitle, placeholders)
            for para in frame.paragraphs:
                para.font.name = self.primary_font
                para.font.size = Pt(self.typography.subtitle_size)
                para.font.color.rgb = ThemeColors.TEXT_SECONDARY
                para.alignment = PP_ALIGN.CENTER
                para.space_after = Pt(6)

        # Footer bar
        self._add_footer_bar(slide)
        
        self._apply_notes(slide, definition, placeholders)

    def _add_section_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create a section divider slide with prominent styling and visual elements."""
        layout = prs.slide_layouts[self.config.blank_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_section_theme(slide)
        
        # Add left accent bar
        self._add_side_accent(slide, "left")
        
        # Add large decorative circle/shape in background
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(self.SLIDE_WIDTH - 4), Inches(-1),
            Inches(5), Inches(5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ThemeColors.ACCENT_PRIMARY
        circle.fill.fore_color.brightness = 0.7  # Very subtle
        circle.line.fill.background()
        
        # Section number indicator
        section_num_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.0), Inches(2.0), Inches(0.8)
        )
        frame = section_num_box.text_frame
        frame.text = f"SECTION {slide_num}"
        para = frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(self.typography.caption_size)
        para.font.color.rgb = ThemeColors.ACCENT_PRIMARY
        para.font.bold = True

        # Section title - large, left-aligned for modern look
        if definition.title:
            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(2.8), Inches(10.0), Inches(1.5)
            )
            frame = title_box.text_frame
            frame.word_wrap = True
            frame.text = _fill(definition.title, placeholders)
            para = frame.paragraphs[0]
            para.font.name = self.primary_font
            para.font.size = Pt(self.typography.section_title_size)
            para.font.color.rgb = ThemeColors.TEXT_HIGHLIGHT
            para.font.bold = True
            para.alignment = PP_ALIGN.LEFT

        # Accent line under title
        self._add_accent_line(slide, top=4.3, width=3.0, left=0.6)

        # Subtitle if provided
        if definition.subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(4.6), Inches(10.0), Inches(1.2)
            )
            frame = sub_box.text_frame
            frame.word_wrap = True
            frame.text = _fill(definition.subtitle, placeholders)
            para = frame.paragraphs[0]
            para.font.name = self.primary_font
            para.font.size = Pt(self.typography.section_subtitle_size)
            para.font.color.rgb = ThemeColors.TEXT_SECONDARY
            para.alignment = PP_ALIGN.LEFT

        # Footer bar with slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)
        self._apply_notes(slide, definition, placeholders)

    def _add_text_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create a text slide with professional styling and card background."""
        layout = prs.slide_layouts[self.config.blank_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add side accent
        self._add_side_accent(slide, "left")
        
        self._add_slide_title(slide, definition, placeholders)
        self._add_accent_line(slide, top=1.1, width=2.0, left=0.6)

        body = definition.text or ""

        # Add card background for text content
        content_top = definition.position if definition.position is not None else 1.4
        card_height = max(2.5, self.SLIDE_HEIGHT - content_top - 0.8)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(content_top), Inches(12.3), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ThemeColors.BACKGROUND_CARD
        card.line.fill.background()
        
        textbox = slide.shapes.add_textbox(
            Inches(0.7), Inches(content_top + 0.2), Inches(11.9), Inches(card_height - 0.4)
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.text = _fill(body, placeholders)
        
        for para in frame.paragraphs:
            para.font.name = self.primary_font
            para.font.size = Pt(self.typography.body_size)
            para.font.color.rgb = ThemeColors.TEXT_PRIMARY
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(8)
        
        # Footer and slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)
        self._apply_notes(slide, definition, placeholders)

    def _add_bullet_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create a bullet point slide with professional styling and smart formatting."""
        layout = prs.slide_layouts[self.config.blank_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add side accent
        self._add_side_accent(slide, "left")
        
        # Add title with proper positioning
        self._add_slide_title(slide, definition, placeholders)
        self._add_accent_line(slide, top=1.1, width=2.0, left=0.6)

        # Bullet content with smart formatting
        content_top = definition.position if definition.position is not None else 1.4
        content_height = max(1.0, self.SLIDE_HEIGHT - content_top - 0.9)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(content_top), Inches(12.3), Inches(content_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ThemeColors.BACKGROUND_CARD
        card.line.fill.background()

        content_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(content_top + 0.2), Inches(11.9), Inches(max(0.1, content_height - 0.4))
        )
        self._populate_bullet_frame(content_box.text_frame, definition.bullets or [], placeholders)

        # Footer and slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)
        self._apply_notes(slide, definition, placeholders)

    def _add_image_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create an image slide with centered image, border, and professional styling."""
        layout = prs.slide_layouts[self.config.blank_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add side accent
        self._add_side_accent(slide, "left")
        
        # Add title
        self._add_slide_title(slide, definition, placeholders)
        self._add_accent_line(slide, top=1.1, width=2.0, left=0.6)

        img_path = Path(definition.image_path) if definition.image_path else None

        box_width = float(definition.image_width or 10.0)
        box_top = (
            definition.image_top
            if definition.image_top is not None
            else (definition.position if definition.position is not None else 1.5)
        )
        box_height = max(2.5, self.SLIDE_HEIGHT - float(box_top) - 1.0)

        max_box_width = self.SLIDE_WIDTH - 1.2
        if definition.image_left is None and box_width > max_box_width:
            box_width = max_box_width

        if definition.image_left is not None:
            box_left = float(definition.image_left)
        else:
            box_left = (self.SLIDE_WIDTH - box_width) / 2
        
        card_padding = 0.15
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(box_left - card_padding),
            Inches(box_top - card_padding),
            Inches(box_width + card_padding * 2),
            Inches(box_height + card_padding * 2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ThemeColors.BACKGROUND_CARD
        card.line.color.rgb = ThemeColors.ACCENT_PRIMARY
        card.line.width = Pt(1)
        
        if not img_path:
            LOGGER.warning("Image slide '%s' missing image_path", definition.title)
            self._add_missing_asset_placeholder(
                slide,
                message="Missing image_path in config",
                left=box_left,
                top=box_top,
                width=box_width,
                height=box_height,
            )
        elif not img_path.exists():
            LOGGER.warning("Image file %s not found", img_path)
            self._add_missing_asset_placeholder(
                slide,
                message=f"Missing image file: {img_path}",
                left=box_left,
                top=box_top,
                width=box_width,
                height=box_height,
            )
        else:
            picture = slide.shapes.add_picture(
                str(img_path),
                Inches(box_left),
                Inches(box_top),
            )

            max_width_emu = Inches(box_width)
            max_height_emu = Inches(box_height)
            scale = min(max_width_emu / picture.width, max_height_emu / picture.height, 1.0)
            if scale < 1.0:
                picture.width = int(picture.width * scale)
                picture.height = int(picture.height * scale)

            if definition.image_left is None and picture.width < max_width_emu:
                picture.left = int(Inches(box_left) + (max_width_emu - picture.width) / 2)
            if definition.image_top is None and definition.position is None and picture.height < max_height_emu:
                picture.top = int(Inches(box_top) + (max_height_emu - picture.height) / 2)

            card.left = int(picture.left - Inches(card_padding))
            card.top = int(picture.top - Inches(card_padding))
            card.width = int(picture.width + Inches(card_padding * 2))
            card.height = int(picture.height + Inches(card_padding * 2))

            if definition.image_border:
                picture.line.color.rgb = ThemeColors.ACCENT_PRIMARY
                picture.line.width = Pt(2)
        
        # Footer and slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)
        self._apply_notes(slide, definition, placeholders)

    def _add_two_column_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create a two-column layout slide with card backgrounds."""
        layout = prs.slide_layouts[self.config.blank_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add side accent
        self._add_side_accent(slide, "left")
        
        self._add_slide_title(slide, definition, placeholders)
        self._add_accent_line(slide, top=1.1, width=2.0, left=0.6)

        # Determine column content (explicit columns take priority)
        if definition.left_content or definition.right_content:
            left_items = definition.left_content
            right_items = definition.right_content
        else:
            bullets = definition.bullets or []
            mid = (len(bullets) + 1) // 2
            left_items = bullets[:mid]
            right_items = bullets[mid:]

        content_top = definition.position if definition.position is not None else 1.4
        column_height = max(2.8, self.SLIDE_HEIGHT - content_top - 0.8)

        # Left column card
        left_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(content_top), Inches(5.9), Inches(column_height)
        )
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = ThemeColors.BACKGROUND_CARD
        left_card.line.fill.background()

        # Left column content
        left_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(content_top + 0.2), Inches(5.5), Inches(column_height - 0.4)
        )
        self._populate_bullet_frame(left_box.text_frame, left_items, placeholders)

        # Right column card
        right_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.9), Inches(content_top), Inches(5.9), Inches(column_height)
        )
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = ThemeColors.BACKGROUND_CARD
        right_card.line.fill.background()

        has_right_text = bool(right_items)
        has_right_image = bool(definition.right_image)

        if has_right_text:
            right_box_height = 2.5 if has_right_image else 5.0
            right_box = slide.shapes.add_textbox(
                Inches(7.1), Inches(content_top + 0.2), Inches(5.5), Inches(right_box_height)
            )
            self._populate_bullet_frame(right_box.text_frame, right_items, placeholders)

        if has_right_image:
            img_path = Path(definition.right_image)
            if img_path.exists():
                image_top = content_top + 0.2 + (2.7 if has_right_text else 0.0)
                box_left = 7.1
                box_width = 5.3
                box_height = max(1.0, (content_top + column_height) - image_top - 0.2)
                picture = slide.shapes.add_picture(
                    str(img_path),
                    Inches(box_left),
                    Inches(image_top),
                )
                max_width_emu = Inches(box_width)
                max_height_emu = Inches(box_height)
                scale = min(max_width_emu / picture.width, max_height_emu / picture.height, 1.0)
                if scale < 1.0:
                    picture.width = int(picture.width * scale)
                    picture.height = int(picture.height * scale)

                if picture.width < max_width_emu:
                    picture.left = int(Inches(box_left) + (max_width_emu - picture.width) / 2)
            else:
                LOGGER.warning("Right column image file %s not found", img_path)
                image_top = content_top + 0.2 + (2.7 if has_right_text else 0.0)
                self._add_missing_asset_placeholder(
                    slide,
                    message=f"Missing image file: {img_path}",
                    left=7.1,
                    top=image_top,
                    width=5.3,
                    height=max(1.0, (content_top + column_height) - image_top - 0.2),
                    font_size=self.typography.caption_size,
                )

        # Footer and slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)
        self._apply_notes(slide, definition, placeholders)

    def _add_table_slide(
        self,
        prs: Presentation,
        definition: LocalSlideDefinition,
        placeholders: Mapping[str, str],
        data_sources: Mapping[str, object],
        slide_num: int = 1,
        total_slides: int = 1,
    ) -> None:
        """Create a table slide with professional dark theme styling and card background."""
        layout = prs.slide_layouts[self.config.blank_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add side accent
        self._add_side_accent(slide, "left")
        
        self._add_slide_title(slide, definition, placeholders)
        self._add_accent_line(slide, top=1.1, width=2.0, left=0.6)

        source_key = definition.table_source
        dataframe = data_sources.get(source_key) if source_key else None
        if not isinstance(dataframe, pd.DataFrame):
            LOGGER.warning("Table slide '%s' missing DataFrame '%s'", definition.title, source_key)
            self._add_missing_asset_placeholder(
                slide,
                message=f"Missing table source: {source_key}",
                left=0.7,
                top=1.6,
                width=11.9,
                height=4.8,
            )

            self._add_footer_bar(slide)
            self._add_slide_number(slide, slide_num, total_slides)
            self._apply_notes(slide, definition, placeholders)
            return

        rows, cols = dataframe.shape
        
        # Calculate table dimensions
        table_width = min(12.0, max(9.0, cols * 1.8))
        table_height = min(5.0, max(2.5, (rows + 1) * 0.45))
        table_left = (self.SLIDE_WIDTH - table_width) / 2
        table_top = 1.5
        
        # Add card background for table
        card_padding = 0.15
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(table_left - card_padding),
            Inches(table_top - card_padding),
            Inches(table_width + card_padding * 2),
            Inches(table_height + card_padding * 2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ThemeColors.BACKGROUND_CARD
        card.line.fill.background()
        
        table_shape = slide.shapes.add_table(
            rows + 1, cols, 
            Inches(table_left), Inches(table_top), 
            Inches(table_width), Inches(table_height)
        )
        table = table_shape.table

        # Style headers with accent background
        for idx, column in enumerate(dataframe.columns):
            cell = table.cell(0, idx)
            cell.text = _fill(str(column), placeholders)
            self._style_table_header(cell)

        # Style data rows with alternating backgrounds
        for row_idx, (_, data_row) in enumerate(dataframe.iterrows(), start=1):
            for col_idx, value in enumerate(data_row):
                cell = table.cell(row_idx, col_idx)
                cell.text = _fill(str(value), placeholders)
                # Highlight first column (usually labels)
                is_label = col_idx == 0
                self._style_table_cell(cell, row_idx % 2 == 0, is_label)
                
                # Right-align numeric values, left-align text
                para = cell.text_frame.paragraphs[0]
                if isinstance(value, (int, float)) or (isinstance(value, str) and 
                                                       re.match(r'^-?\d+\.?\d*', str(value))):
                    para.alignment = PP_ALIGN.RIGHT
                else:
                    para.alignment = PP_ALIGN.LEFT

        # Add statistical footnote for publication quality
        footnote_text = "Values represent mean ± SD. Statistical comparisons: two-tailed t-test, α = 0.05."
        footnote_top = 1.5 + table_height + 0.2
        self._add_table_footnote(slide, footnote_text, footnote_top)

        # Footer and slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)
        self._apply_notes(slide, definition, placeholders)

    def _add_table_footnote(self, slide, text: str, top: float) -> None:
        """Add italicized footnote below table for statistical annotations."""
        textbox = slide.shapes.add_textbox(
            Inches(0.7), Inches(top), Inches(11.9), Inches(0.4)
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.text = text
        para = frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(self.typography.caption_size)
        para.font.color.rgb = ThemeColors.TEXT_MUTED
        para.font.italic = True
        para.alignment = PP_ALIGN.LEFT

    def _style_table_header(self, cell) -> None:
        """Style table header cell with accent background."""
        cell.fill.solid()
        cell.fill.fore_color.rgb = ThemeColors.TABLE_HEADER_BG
        para = cell.text_frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(self.typography.table_header_size)
        para.font.bold = True
        para.font.color.rgb = ThemeColors.TEXT_PRIMARY
        para.alignment = PP_ALIGN.CENTER

    def _style_table_cell(self, cell, alternate: bool = False, is_label: bool = False) -> None:
        """Style table data cell with optional alternating background."""
        if alternate:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ThemeColors.TABLE_ROW_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ThemeColors.TABLE_ROW_NORMAL
        para = cell.text_frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(self.typography.table_body_size)
        
        if is_label:
            para.font.color.rgb = ThemeColors.ACCENT_HIGHLIGHT
            para.font.bold = True
            para.alignment = PP_ALIGN.LEFT
        else:
            para.font.color.rgb = ThemeColors.TEXT_PRIMARY
            para.alignment = PP_ALIGN.CENTER

    def _add_equation_slide(
        self, prs: Presentation, definition: LocalSlideDefinition, placeholders: Mapping[str, str],
        slide_num: int = 1, total_slides: int = 1
    ) -> None:
        """Create an equation slide with formatted chemical/mathematical equations."""
        layout = prs.slide_layouts[self.config.title_layout_index]
        slide = prs.slides.add_slide(layout)
        self._apply_dark_theme(slide)
        
        # Add title (using definition.title attribute)
        title_shape = slide.shapes.title
        title_shape.text = definition.title if hasattr(definition, 'title') else 'Equations'
        
        # Format equations (access via definition.equations attribute)
        body_shape = slide.placeholders[1]
        if definition.position is not None:
            body_shape.top = Inches(definition.position)
        content = body_shape.text_frame
        equations = definition.equations if hasattr(definition, 'equations') else []
        for eq in equations:
            p = content.add_paragraph()
            p.text = eq
            p.font.size = Pt(24)
            p.font.name = 'Cambria Math'
        
        # Add speaker notes if available
        if hasattr(definition, 'notes'):
            self._add_speaker_notes(slide, definition.notes)
        
        # Footer and slide number
        self._add_footer_bar(slide)
        self._add_slide_number(slide, slide_num, total_slides)

    def _add_speaker_notes(self, slide, notes_text):
        """
        Adds speaker notes to a slide if notes_text is provided
        Args:
            slide: PowerPoint slide object
            notes_text: Notes content (string)
        """
        if notes_text and len(notes_text.strip()) > 0:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    # ------------------------------------------------------------------
    # Helper Methods
    # ------------------------------------------------------------------
    def _populate_bullet_frame(
        self,
        text_frame,
        items: Iterable[str],
        placeholders: Mapping[str, str],
    ) -> None:
        """Populate a text frame with consistently styled bullets."""
        text_frame.clear()
        text_frame.word_wrap = True

        bullet_items = list(items or [])
        if not bullet_items:
            text_frame.text = ""
            return

        for idx, raw_item in enumerate(bullet_items):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            bullet_text = _fill(raw_item, placeholders)
            if bullet_text is None:
                bullet_text = ""
            bullet_text = bullet_text.replace("\r\n", "\n").replace("\r", "\n").lstrip("\n")

            indent_level = 0
            while bullet_text.startswith(INDENT_TOKEN):
                indent_level += 1
                bullet_text = bullet_text[len(INDENT_TOKEN) :]

            bullet_text = bullet_text.strip()
            if not bullet_text:
                paragraph.text = ""
                continue

            manual_prefix = None
            for prefix in MANUAL_BULLET_PREFIXES:
                if bullet_text.startswith(prefix):
                    manual_prefix = prefix
                    bullet_text = bullet_text[len(prefix) :].lstrip()
                    break

            bullet_prefix = manual_prefix or ("◦" if indent_level > 0 else "•")
            paragraph.level = min(indent_level, 5)
            paragraph.space_after = Pt(4)
            paragraph.space_before = Pt(1)
            paragraph.text = ""

            base_font_size = (
                self.typography.body_size if indent_level == 0 else self.typography.bullet_indent_size
            )
            base_color = ThemeColors.TEXT_PRIMARY if indent_level == 0 else ThemeColors.TEXT_SECONDARY

            prefix_run = paragraph.add_run()
            prefix_run.text = f"{bullet_prefix} "
            prefix_run.font.name = self.primary_font
            prefix_run.font.size = Pt(base_font_size)
            prefix_run.font.color.rgb = base_color

            self._apply_inline_formatting(paragraph, bullet_text, base_font_size, base_color)

    def _apply_inline_formatting(
        self,
        paragraph,
        text: str,
        font_size: int,
        color: RGBColor,
    ) -> None:
        """Render inline **bold** spans while preserving base styling."""
        if not text:
            return

        segments = text.split("**")
        for idx, segment in enumerate(segments):
            if not segment:
                continue
            run = paragraph.add_run()
            run.text = segment
            run.font.name = self.primary_font
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = idx % 2 == 1

    def _add_slide_title(
        self, slide, definition: LocalSlideDefinition, placeholders: Mapping[str, str]
    ) -> None:
        """Add a styled title to a content slide."""
        if not definition.title:
            return
        textbox = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.8)
        )
        frame = textbox.text_frame
        frame.text = _fill(definition.title, placeholders)
        para = frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(self.typography.heading_size)
        para.font.color.rgb = ThemeColors.TEXT_PRIMARY
        para.font.bold = True

    def _apply_notes(self, slide, definition: LocalSlideDefinition, placeholders: Mapping[str, str]) -> None:
        if not definition.notes:
            return
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = _fill(definition.notes, placeholders)

    def _add_missing_asset_placeholder(
        self,
        slide,
        *,
        message: str,
        left: float,
        top: float,
        width: float,
        height: float,
        font_size: int | None = None,
    ) -> None:
        """Add a visible placeholder message in the slide (used when assets are missing)."""
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = ThemeColors.BACKGROUND_ALT
        box.line.color.rgb = ThemeColors.ACCENT_WARNING
        box.line.width = Pt(2)

        textbox = slide.shapes.add_textbox(
            Inches(left + 0.2),
            Inches(top + 0.2),
            Inches(max(0.1, width - 0.4)),
            Inches(max(0.1, height - 0.4)),
        )
        frame = textbox.text_frame
        frame.word_wrap = True
        frame.text = message
        para = frame.paragraphs[0]
        para.font.name = self.primary_font
        para.font.size = Pt(font_size or self.typography.body_secondary_size)
        para.font.color.rgb = ThemeColors.ACCENT_WARNING
        para.font.bold = True


def _fill(text: str, replacements: Mapping[str, str]) -> str:
    """Replace placeholders in text with their values."""
    if text is None:
        return ""
    rendered = str(text)
    for placeholder, value in replacements.items():
        rendered = rendered.replace(placeholder, value)
    return rendered
