from __future__ import annotations

import json
import logging
import re
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import asdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, Iterable, List, Mapping, Optional

import pandas as pd
from pptx import Presentation

from .config import ProjectConfig

LOGGER = logging.getLogger(__name__)
EMU_PER_IN = 914400

_PLACEHOLDER_PATTERN = re.compile(r"\{[A-Z0-9_]+\}")

_PPTX_XML_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
_PPTX_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


# =============================================================================
# Scientific Formatting Helpers
# =============================================================================

def _format_pvalue(p: float) -> str:
    """Format p-value following scientific convention.
    
    Args:
        p: P-value from statistical test
        
    Returns:
        Formatted string: 'p < 0.001' or 'p = 0.042'
    """
    if p < 0.001:
        return "p < 0.001"
    elif p < 0.01:
        return f"p = {p:.3f}"
    else:
        return f"p = {p:.2f}"


def _format_measurement(value: float, uncertainty: float = None, unit: str = "") -> str:
    """Format measurement with optional uncertainty.
    
    Args:
        value: Measured value
        uncertainty: Standard error or uncertainty (optional)
        unit: Unit string (optional)
        
    Returns:
        Formatted string: '10.5 ± 0.2 nm' or '10.5 nm'
    """
    if uncertainty is None:
        return f"{value:.3g} {unit}".strip()
    return f"{value:.3g} ± {uncertainty:.2g} {unit}".strip()


def _format_ci(low: float, high: float, unit: str = "", level: int = 95) -> str:
    """Format confidence interval.
    
    Args:
        low: Lower bound
        high: Upper bound
        unit: Unit string (optional)
        level: Confidence level percentage (default 95)
        
    Returns:
        Formatted string: '95% CI: [2.1, 3.4] ppm'
    """
    return f"{level}% CI: [{low:.3g}, {high:.3g}] {unit}".strip()


def _as_utc_iso(ts: Optional[datetime] = None) -> str:
    dt = ts or datetime.now(timezone.utc)
    return dt.astimezone(timezone.utc).isoformat(timespec="seconds")


def apply_auto_metrics(
    *,
    placeholders: Dict[str, str],
    data_sources: Mapping[str, object],
) -> Dict[str, str]:
    """Derive key placeholders from loaded data sources.

    This reduces drift between tables/figures and the text in placeholders.
    Values are only applied when the corresponding placeholder is missing.
    """

    updated = dict(placeholders)

    def _set_if_missing(key: str, value: Any) -> None:
        if key in updated and str(updated.get(key, "")).strip():
            return
        if value is None:
            return
        value_str = str(value).strip()
        if not value_str:
            return
        updated[key] = value_str

    def _set_if_blank_in_template(key: str, value: Any) -> None:
        if str(placeholders.get(key, "")).strip():
            return
        if value is None:
            return
        value_str = str(value).strip()
        if not value_str:
            return
        updated[key] = value_str

    comparison = data_sources.get("comparison_table")
    if isinstance(comparison, pd.DataFrame) and {"Metric", "Reference Paper", "This Work"}.issubset(
        set(comparison.columns)
    ):
        by_metric = {
            str(row["Metric"]).strip(): row
            for _, row in comparison.iterrows()
            if str(row.get("Metric", "")).strip()
        }

        roi_row = by_metric.get("ROI")
        if roi_row is not None:
            _set_if_missing("{LITERATURE_ROI}", roi_row.get("Reference Paper"))
            _set_if_missing("{OPTIMAL_ROI}", roi_row.get("This Work"))

        lod_row = by_metric.get("LoD")
        lod_paper = None
        lod_this = None
        if lod_row is not None:
            lod_paper = lod_row.get("Reference Paper")
            lod_this = lod_row.get("This Work")
            _set_if_missing("{LOD_PAPER}", lod_paper)
            _set_if_missing("{LOD_THIS_WORK}", lod_this)

        r2_row = by_metric.get("R²")
        if r2_row is not None:
            _set_if_missing("{R2}", r2_row.get("This Work"))

        loocv_row = by_metric.get("LOOCV R²")
        if loocv_row is not None:
            _set_if_missing("{LOOCV_R2}", loocv_row.get("This Work"))

        t90_row = by_metric.get("T₉₀")
        if t90_row is not None:
            _set_if_missing("{T90}", t90_row.get("This Work"))

        rec_row = by_metric.get("Recovery")
        if rec_row is not None:
            _set_if_missing("{RECOVERY}", rec_row.get("This Work"))

        sens_row = by_metric.get("Sensitivity")
        if sens_row is not None:
            _set_if_missing("{SENSITIVITY}", sens_row.get("This Work"))
        try:
            paper_val = float(str(lod_paper).split()[0]) if lod_paper is not None else None
            this_val = float(str(lod_this).split()[0]) if lod_this is not None else None
            if paper_val and this_val and this_val > 0:
                factor = paper_val / this_val
                factor_rounded = int(round(factor))
                _set_if_missing("{IMPROVEMENT_FACTOR}", f"{factor_rounded}×")
        except Exception:
            pass

    voc_protocol = data_sources.get("voc_protocol")
    if isinstance(voc_protocol, pd.DataFrame) and "Files" in voc_protocol.columns:
        files_vals = pd.to_numeric(voc_protocol["Files"], errors="coerce").dropna().tolist()
        if files_vals:
            total = int(sum(files_vals))
            _set_if_missing("{SAMPLE_SIZE}", f"n = {total:,} spectra")

    reference_metrics = data_sources.get("reference_metrics")
    if isinstance(reference_metrics, pd.DataFrame) and {"Metric", "Value"}.issubset(
        set(reference_metrics.columns)
    ):
        metric_map = {
            str(row["Metric"]).strip(): str(row.get("Value", "")).strip()
            for _, row in reference_metrics.iterrows()
            if str(row.get("Metric", "")).strip()
        }
        _set_if_missing("{ZNO_THICKNESS}", metric_map.get("ZnO Thickness"))

    acetone_metrics = data_sources.get("acetone_calibration_metrics")
    if isinstance(acetone_metrics, dict):
        centroid = (
            acetone_metrics.get("calibration_wavelength_shift", {})
            .get("centroid", {})
        )
        p_value = centroid.get("p_value")
        if isinstance(p_value, (int, float)):
            _set_if_blank_in_template("{P_VALUE}", _format_pvalue(p_value))

        slope_ci = centroid.get("slope_ci_95")
        if isinstance(slope_ci, list) and len(slope_ci) == 2:
            try:
                lo = float(slope_ci[0])
                hi = float(slope_ci[1])
                _set_if_blank_in_template("{CI_95}", _format_ci(lo, hi, "nm/ppm"))
            except Exception:
                pass

        noise_std = centroid.get("noise_std")
        if isinstance(noise_std, (int, float)):
            _set_if_blank_in_template("{RESIDUAL_STD}", f"{noise_std:.4f} nm")

        r2_val = centroid.get("r2")
        if isinstance(r2_val, (int, float)):
            _set_if_blank_in_template("{R2}", f"{r2_val:.4f}")

        slope_val = centroid.get("slope")
        unit = centroid.get("slope_unit")
        if isinstance(slope_val, (int, float)):
            unit_str = str(unit).strip() if unit else "nm/ppm"
            _set_if_blank_in_template("{SENSITIVITY}", f"{slope_val:.4f} {unit_str}")

        lod_val = centroid.get("lod_ppm")
        if isinstance(lod_val, (int, float)):
            _set_if_blank_in_template("{LOD_THIS_WORK}", f"{lod_val:.2f} ppm")

        r2_cv = centroid.get("r2_cv")
        if isinstance(r2_cv, (int, float)):
            _set_if_blank_in_template("{LOOCV_R2}", f"{r2_cv:.4f}")

    return updated


def collect_referenced_paths(config: ProjectConfig) -> List[Path]:
    paths: List[Path] = []

    if config.pptx and config.pptx.template_path:
        paths.append(Path(config.pptx.template_path))

    for slide in config.slides:
        if slide.image_path:
            paths.append(Path(slide.image_path))
        if slide.right_image:
            paths.append(Path(slide.right_image))

    if config.google_slides:
        for upload in config.google_slides.image_uploads:
            paths.append(Path(upload.path))

    for _, location in (config.data_sources or {}).items():
        paths.append(Path(location))

    return paths


def check_placeholder_consistency(
    *,
    placeholders: Mapping[str, str],
    data_sources: Mapping[str, object],
) -> List[str]:
    """Warn if key headline placeholders disagree with loaded metric tables."""

    warnings: List[str] = []

    def _split_numeric_unit(value: str) -> tuple[list[float], str]:
        raw = str(value)
        numbers = re.findall(r"[-+]?\d*\.?\d+(?:[eE][-+]?\d+)?", raw)
        floats: list[float] = []
        for num in numbers:
            try:
                floats.append(float(num))
            except ValueError:
                continue
        unit = re.sub(r"[-+]?\d*\.?\d+(?:[eE][-+]?\d+)?", "", raw)
        unit = re.sub(r"\s+", " ", unit).strip()
        return floats, unit

    def _normalize_unit(unit: str) -> str:
        raw = str(unit)
        raw = raw.replace("±", " ")
        raw = raw.replace("≈", " ")
        raw = raw.replace("~", " ")
        raw = raw.replace("–", " ")
        raw = raw.replace("-", " ")
        raw = raw.replace("+/-", " ")
        raw = re.sub(r"[(),\[\]]", " ", raw)
        raw = re.sub(r"\s+", " ", raw).strip()
        return raw

    def _looks_like_uncertainty(raw: str) -> bool:
        collapsed = str(raw).replace(" ", "")
        return ("±" in collapsed) or ("+/-" in collapsed)

    comparison = data_sources.get("comparison_table")
    if not isinstance(comparison, pd.DataFrame) or "Metric" not in comparison.columns:
        return warnings

    by_metric = {
        str(row["Metric"]).strip(): row
        for _, row in comparison.iterrows()
        if str(row.get("Metric", "")).strip()
    }

    expected_map = {
        "{LOD_PAPER}": ("LoD", "Reference Paper"),
        "{LOD_THIS_WORK}": ("LoD", "This Work"),
        "{R2}": ("R²", "This Work"),
        "{LOOCV_R2}": ("LOOCV R²", "This Work"),
        "{SENSITIVITY}": ("Sensitivity", "This Work"),
        "{OPTIMAL_ROI}": ("ROI", "This Work"),
        "{LITERATURE_ROI}": ("ROI", "Reference Paper"),
        "{T90}": ("T₉₀", "This Work"),
        "{RECOVERY}": ("Recovery", "This Work"),
    }

    for placeholder_key, (metric_key, column) in expected_map.items():
        if placeholder_key not in placeholders:
            continue
        row = by_metric.get(metric_key)
        if row is None:
            continue
        table_value = row.get(column)
        if table_value is None:
            continue
        table_value_str = str(table_value).strip()
        placeholder_value_str = str(placeholders.get(placeholder_key, "")).strip()
        if not table_value_str or not placeholder_value_str:
            continue

        table_nums, table_unit = _split_numeric_unit(table_value_str)
        ph_nums, ph_unit = _split_numeric_unit(placeholder_value_str)

        table_unit_n = _normalize_unit(table_unit)
        ph_unit_n = _normalize_unit(ph_unit)

        # If both look like a single numeric value with matching unit, compare with tolerance.
        if table_nums and ph_nums and table_unit_n == ph_unit_n:
            tol = 1e-3
            if abs(table_nums[0] - ph_nums[0]) <= tol:
                if (
                    len(table_nums) == 2
                    and len(ph_nums) == 2
                    and not _looks_like_uncertainty(table_value_str)
                    and not _looks_like_uncertainty(placeholder_value_str)
                ):
                    if abs(table_nums[1] - ph_nums[1]) <= tol:
                        continue
                else:
                    continue

        if table_value_str != placeholder_value_str:
            warnings.append(
                f"⚠️  Placeholder {placeholder_key}='{placeholder_value_str}' differs from "
                f"comparison_table[{metric_key}][{column}]='{table_value_str}'. "
                f"Consider using auto_metrics or updating the placeholder."
            )

    return warnings


def _try_read_text(path: Path) -> Optional[str]:
    try:
        return path.read_text(encoding="utf-8").strip()
    except Exception:
        return None


def get_git_revision(repo_root: Path) -> Optional[str]:
    """Best-effort git revision detection without invoking subprocess."""
    git_dir = repo_root / ".git"
    head_path = git_dir / "HEAD"
    if not head_path.exists():
        return None

    head = _try_read_text(head_path)
    if not head:
        return None

    if head.startswith("ref:"):
        ref = head.split(":", 1)[1].strip()
        ref_path = git_dir / ref
        sha = _try_read_text(ref_path)
        if sha:
            return sha

        packed = _try_read_text(git_dir / "packed-refs")
        if packed:
            for line in packed.splitlines():
                line = line.strip()
                if not line or line.startswith("#") or line.startswith("^"):
                    continue
                try:
                    sha_candidate, ref_candidate = line.split(" ", 1)
                except ValueError:
                    continue
                if ref_candidate.strip() == ref:
                    return sha_candidate.strip()
        return None

    return head


def validate_project(
    *,
    config: ProjectConfig,
    placeholders: Mapping[str, str],
    strict: bool,
    generate_google: bool = True,
    generate_pptx: bool = True,
) -> List[str]:
    """Validate referenced assets and placeholders.

    Returns a list of warning strings. In strict mode, missing assets raise.
    """

    warnings: List[str] = []

    referenced: List[Path] = []

    if generate_pptx:
        if config.pptx and config.pptx.template_path:
            referenced.append(Path(config.pptx.template_path))
        for slide in config.slides:
            if slide.image_path:
                referenced.append(Path(slide.image_path))
            if slide.right_image:
                referenced.append(Path(slide.right_image))
        for _, location in (config.data_sources or {}).items():
            referenced.append(Path(location))

    if generate_google:
        if config.credentials_file:
            referenced.append(Path(config.credentials_file))
        if config.google_slides:
            for upload in config.google_slides.image_uploads:
                referenced.append(Path(upload.path))

    missing_paths = [p for p in referenced if not p.exists()]
    if missing_paths:
        msg = "Missing referenced files:\n" + "\n".join(f"- {p}" for p in missing_paths)
        if strict:
            raise FileNotFoundError(msg)
        warnings.append(msg)

    unresolved = find_unresolved_placeholders_in_slide_defs(config, placeholders)
    if unresolved:
        msg = "Unresolved placeholders found in slide definitions:\n" + "\n".join(
            f"- {item}" for item in sorted(unresolved)
        )
        if strict:
            raise ValueError(msg)
        warnings.append(msg)

    return warnings


def find_unresolved_placeholders_in_slide_defs(
    config: ProjectConfig, placeholders: Mapping[str, str]
) -> List[str]:
    known = set(placeholders.keys())

    def _scan(text: Optional[str], context: str) -> Iterable[str]:
        if not text:
            return []
        hits = sorted(set(_PLACEHOLDER_PATTERN.findall(str(text))))
        return [f"{context}: {hit}" for hit in hits if hit not in known]

    issues: List[str] = []
    for idx, slide in enumerate(config.slides, start=1):
        issues.extend(_scan(slide.title, f"slide[{idx}].title"))
        issues.extend(_scan(slide.subtitle, f"slide[{idx}].subtitle"))
        issues.extend(_scan(slide.text, f"slide[{idx}].text"))
        issues.extend(_scan(slide.notes, f"slide[{idx}].notes"))
        for b_idx, bullet in enumerate(slide.bullets or [], start=1):
            issues.extend(_scan(bullet, f"slide[{idx}].bullets[{b_idx}]"))
        for l_idx, item in enumerate(slide.left_content or [], start=1):
            issues.extend(_scan(item, f"slide[{idx}].left_content[{l_idx}]"))
        for r_idx, item in enumerate(slide.right_content or [], start=1):
            issues.extend(_scan(item, f"slide[{idx}].right_content[{r_idx}]"))
        for e_idx, eq in enumerate(slide.equations or [], start=1):
            issues.extend(_scan(eq, f"slide[{idx}].equations[{e_idx}]"))

    return issues


def scan_pptx_for_placeholders(pptx_path: str | Path) -> List[str]:
    prs = Presentation(str(pptx_path))
    hits: List[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if not text:
                continue
            for match in sorted(set(_PLACEHOLDER_PATTERN.findall(text))):
                hits.append(f"slide[{slide_idx}]: {match}")
    return hits


def scan_pptx_for_picture_margin_issues_xml(
    *,
    pptx_path: str | Path,
    min_margin_in: float = 0.25,
) -> tuple[Dict[int, Dict[str, int]], Dict[str, int]]:
    per_slide: Dict[int, Dict[str, int]] = {}
    meta: Dict[str, int] = {}

    margin_emu = int(float(min_margin_in) * EMU_PER_IN)

    def q(tag: str) -> str:
        prefix, local = tag.split(":", 1)
        return f"{{{_PPTX_XML_NS[prefix]}}}{local}"

    pptx_path = Path(pptx_path)
    with zipfile.ZipFile(pptx_path, "r") as z:
        pres = ET.fromstring(z.read("ppt/presentation.xml"))
        sld_sz = pres.find(q("p:sldSz"))
        if sld_sz is None:
            raise ValueError("Could not find slide size in ppt/presentation.xml")
        slide_w = int(sld_sz.attrib.get("cx", "0"))
        slide_h = int(sld_sz.attrib.get("cy", "0"))
        meta.update({"slide_w": slide_w, "slide_h": slide_h, "margin_emu": margin_emu})

        slide_id_list = pres.find(q("p:sldIdLst"))
        rels = ET.fromstring(z.read("ppt/_rels/presentation.xml.rels"))
        rel_map = {
            rel.attrib.get("Id"): rel.attrib.get("Target")
            for rel in rels.findall(f"{{{_PPTX_REL_NS}}}Relationship")
        }

        ordered_slide_paths: list[str] = []
        if slide_id_list is not None:
            for slide_id in slide_id_list.findall(q("p:sldId")):
                rid = slide_id.attrib.get(q("r:id"))
                target = rel_map.get(rid)
                if not target:
                    continue
                ordered_slide_paths.append("ppt/" + target.lstrip("/"))

        if not ordered_slide_paths:
            ordered_slide_paths = sorted(
                [
                    name
                    for name in z.namelist()
                    if name.startswith("ppt/slides/slide") and name.endswith(".xml")
                ]
            )

        for slide_idx, slide_xml_path in enumerate(ordered_slide_paths, start=1):
            root = ET.fromstring(z.read(slide_xml_path))

            overflow = 0
            tight = 0

            for pic in root.findall(".//" + q("p:pic")):
                xfrm = pic.find(".//" + q("a:xfrm"))
                if xfrm is None:
                    continue
                off = xfrm.find(q("a:off"))
                ext = xfrm.find(q("a:ext"))
                if off is None or ext is None:
                    continue

                x = int(off.attrib.get("x", "0"))
                y = int(off.attrib.get("y", "0"))
                cx = int(ext.attrib.get("cx", "0"))
                cy = int(ext.attrib.get("cy", "0"))

                r = x + cx
                b = y + cy

                is_overflow = (x < 0) or (y < 0) or (r > slide_w) or (b > slide_h)
                is_tight = (
                    (x < margin_emu)
                    or (y < margin_emu)
                    or ((slide_w - r) < margin_emu)
                    or ((slide_h - b) < margin_emu)
                )

                if is_overflow:
                    overflow += 1
                if is_tight:
                    tight += 1

            if overflow or tight:
                per_slide[slide_idx] = {"overflow": overflow, "tight": tight}

    return per_slide, meta


def _compute_file_hash(path: Path) -> str:
    """Compute SHA256 hash of a file for reproducibility tracking.
    
    Args:
        path: Path to file
        
    Returns:
        Hexadecimal SHA256 hash string, or 'ERROR' if computation fails
    """
    import hashlib
    
    sha256 = hashlib.sha256()
    try:
        with path.open('rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                sha256.update(chunk)
        return sha256.hexdigest()
    except Exception:
        return "ERROR"


def write_manifest(
    *,
    output_pptx_path: str | Path,
    config: ProjectConfig,
    title: str,
    placeholders: Mapping[str, str],
    warnings: List[str],
    google_slides_link: Optional[str] = None,
) -> Path:
    out_path = Path(output_pptx_path)
    manifest_path = out_path.with_suffix(out_path.suffix + ".manifest.json")

    config_snapshot_path: Optional[Path] = None
    if config.source_path:
        src = Path(config.source_path)
        if src.exists() and src.is_file():
            config_snapshot_path = out_path.with_suffix(out_path.suffix + src.suffix + ".config")
            try:
                config_snapshot_path.write_text(src.read_text(encoding="utf-8"), encoding="utf-8")
            except Exception:
                config_snapshot_path = None

    repo_root = Path(__file__).resolve().parents[1]
    git_revision = get_git_revision(repo_root)

    referenced_files = []
    for path in collect_referenced_paths(config):
        try:
            stat = path.stat()
            referenced_files.append(
                {
                    "path": str(path),
                    "exists": True,
                    "size_bytes": stat.st_size,
                    "mtime_utc": datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc).isoformat(
                        timespec="seconds"
                    ),
                    "sha256": _compute_file_hash(path),
                }
            )
        except FileNotFoundError:
            referenced_files.append({"path": str(path), "exists": False})

    import sys
    import numpy as np
    
    payload: Dict[str, Any] = {
        "created_utc": _as_utc_iso(),
        "git_revision": git_revision,
        "environment": {
            "python_version": sys.version.split()[0],
            "pandas_version": pd.__version__,
            "numpy_version": np.__version__ if 'numpy' in sys.modules else None,
        },
        "title": title,
        "pptx_path": str(out_path),
        "config_source": config.source_path,
        "config_snapshot": str(config_snapshot_path) if config_snapshot_path else None,
        "project_name": config.project_name,
        "placeholders": dict(placeholders),
        "data_sources": dict(config.data_sources or {}),
        "referenced_files": referenced_files,
        "warnings": list(warnings),
        "google_slides_link": google_slides_link,
        "google_slides": asdict(config.google_slides) if config.google_slides else None,
        "pptx": asdict(config.pptx) if config.pptx else None,
    }

    manifest_path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    return manifest_path


def collect_markdown_image_paths(markdown_text: str) -> List[str]:
    pattern = re.compile(r"!\[[^\]]*\]\(([^\)]+)\)")
    results: List[str] = []
    for match in pattern.finditer(markdown_text):
        raw = match.group(1).strip()
        if not raw:
            continue
        if raw.startswith("<") and raw.endswith(">"):
            raw = raw[1:-1].strip()
        # Drop optional title portion: (path "title")
        if " " in raw:
            raw = raw.split(" ", 1)[0].strip()
        raw = raw.strip("\"")
        if not raw:
            continue
        if raw.lower().startswith("http://") or raw.lower().startswith("https://"):
            continue
        results.append(raw)
    return sorted(set(results))


def scan_markdown_for_missing_assets(
    *, markdown_path: Path, project_root: Optional[Path] = None
) -> List[str]:
    text = markdown_path.read_text(encoding="utf-8")
    base = project_root.resolve() if project_root else markdown_path.parent.resolve()

    warnings: List[str] = []
    for rel in collect_markdown_image_paths(text):
        candidate = Path(rel)
        resolved = candidate if candidate.is_absolute() else (base / candidate)
        if not resolved.exists():
            warnings.append(f"Missing asset referenced in {markdown_path.name}: {rel}")
    return warnings


def write_artifact_bundle(
    *,
    bundle_path: Path,
    project_root: Path,
    files: Iterable[Path],
) -> Path:
    bundle_path.parent.mkdir(parents=True, exist_ok=True)
    root = project_root.resolve()

    unique: List[Path] = []
    seen = set()
    for path in files:
        resolved = Path(path).resolve()
        if resolved in seen:
            continue
        seen.add(resolved)
        unique.append(resolved)

    with zipfile.ZipFile(bundle_path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for file_path in unique:
            if not file_path.exists() or not file_path.is_file():
                continue
            try:
                arcname = str(file_path.relative_to(root))
            except ValueError:
                arcname = file_path.name
            zf.write(file_path, arcname=arcname)

    return bundle_path


def scan_markdown_for_submission_issues(*, markdown_path: Path) -> List[str]:
    text = markdown_path.read_text(encoding="utf-8")
    hits: List[str] = []

    blocker_patterns = [
        (r"\btemporary\b", "Temporary placeholders present (replace authors/affiliations/corresponding author before submission)."),
        (r"\{[A-Z0-9_]+\}", "Unresolved placeholder token present."),
        (r"\[To be added\]", "Draft placeholder '[To be added]' present."),
        (r"\bTBD\b", "Draft placeholder 'TBD' present."),
        (r"corresponding\.author@", "Temporary corresponding author email present."),
        (r"\[Institution\]", "Author affiliation placeholder '[Institution]' present."),
        (r"\[City,\s*Country\]", "Author affiliation placeholder '[City, Country]' present."),
        (r"\[corresponding\.email@institution\]", "Corresponding author email placeholder '[corresponding.email@institution]' present."),
    ]

    advisory_patterns = [
        (
            r"Data and code will be made available upon reasonable request\.",
            "Data availability states 'upon reasonable request' (consider adding a repository link for submission).",
        ),
    ]

    for pat, msg in blocker_patterns:
        if re.search(pat, text, flags=re.IGNORECASE):
            hits.append(f"BLOCKER: {msg}")

    for pat, msg in advisory_patterns:
        if re.search(pat, text, flags=re.IGNORECASE):
            hits.append(f"ADVISORY: {msg}")

    return hits
