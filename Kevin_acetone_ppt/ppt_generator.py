from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


class ResearchPresentationTemplate:
    """
    A professional dark-themed research presentation template for Python-PPTX.
    Optimized for academic and scientific presentations with high contrast and readability.
    """

    # Color scheme - Dark theme with blue accent
    COLORS = {
        'background': RGBColor(15, 23, 42),    # slate-900
        'text_primary': RGBColor(241, 245, 249),  # slate-50
        'text_secondary': RGBColor(203, 213, 225), # slate-300
        'accent': RGBColor(56, 189, 248),      # sky-400
        'accent_dark': RGBColor(14, 165, 233),  # sky-500
    }

    def __init__(self, filename: str = 'presentation.pptx'):
        """
        Initialize the presentation with dark theme styling.

        Args:
            filename: Output filename for the presentation
        """
        self.filename = filename
        self.presentation = Presentation()

        # Set slide size to 16:9 widescreen
        self.presentation.slide_width = Inches(13.33)
        self.presentation.slide_height = Inches(7.5)

        # Apply dark theme to all slides
        self._apply_theme()

    def _apply_theme(self):
        """Apply consistent dark theme across all slides."""
        # This would set global theme if python-pptx supported it fully
        # For now, styling is applied per slide in the methods
        pass

    def _add_background(self, slide):
        """Add dark background to a slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['background']

    def _style_title(self, title_shape, font_size=44):
        """Apply consistent styling to title text."""
        if title_shape:
            title_shape.text_frame.paragraphs[0].font.name = 'Arial'
            title_shape.text_frame.paragraphs[0].font.size = Pt(font_size)
            title_shape.text_frame.paragraphs[0].font.color.rgb = self.COLORS['text_primary']
            title_shape.text_frame.paragraphs[0].font.bold = True

    def _style_text(self, text_frame, font_size=24):
        """Apply consistent styling to body text."""
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = self.COLORS['text_primary']

    def add_title_slide(self, title: str, subtitle: str = None, author: str = None, date: str = None):
        """
        Add a professional title slide.

        Args:
            title: Main presentation title
            subtitle: Optional subtitle
            author: Author name
            date: Presentation date
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        self._add_background(slide)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title(slide.shapes.title, 48)

        # Subtitle and metadata
        if subtitle or author or date:
            content = []
            if subtitle:
                content.append(f"Subtitle: {subtitle}")
            if author:
                content.append(f"Author: {author}")
            if date:
                content.append(f"Date: {date}")

            textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(3))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            for item in content:
                paragraph = text_frame.paragraphs[0] if len(text_frame.paragraphs) == 0 else text_frame.add_paragraph()
                paragraph.text = item
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = self.COLORS['text_secondary']

    def add_content_slide(self, title: str, content_items: list):
        """
        Add a content slide with bullet points.

        Args:
            title: Slide title
            content_items: List of bullet points
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        self._add_background(slide)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title(slide.shapes.title, 36)

        # Content
        if hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            for item in content_items:
                paragraph = text_frame.add_paragraph()
                paragraph.text = item
                paragraph.level = 0
                self._style_text(text_frame)

    def add_section_slide(self, section_title: str):
        """
        Add a section divider slide.

        Args:
            section_title: Section title (e.g., "Methodology", "Results")
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        self._add_background(slide)

        # Centered section title
        textbox = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9), Inches(2))
        text_frame = textbox.text_frame
        text_frame.text = section_title
        text_frame.paragraphs[0].font.name = 'Arial'
        text_frame.paragraphs[0].font.size = Pt(48)
        text_frame.paragraphs[0].font.color.rgb = self.COLORS['accent']
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_two_column_slide(self, title: str, left_content: list, right_content: list):
        """
        Add a two-column content slide.

        Args:
            title: Slide title
            left_content: List of items for left column
            right_content: List of items for right column
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[1])
        self._add_background(slide)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_title(slide.shapes.title, 36)

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True

        for item in left_content:
            paragraph = left_frame.paragraphs[0] if len(left_frame.paragraphs) == 0 else left_frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            self._style_text(left_frame, 22)

        # Right column
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True

        for item in right_content:
            paragraph = right_frame.paragraphs[0] if len(right_frame.paragraphs) == 0 else right_frame.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            self._style_text(right_frame, 22)

    def add_closing_slide(self, title: str = "Thank You!", subtitle: str = None, contact: str = None):
        """
        Add a closing slide.

        Args:
            title: Closing title
            subtitle: Optional subtitle
            contact: Contact information
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[0])
        self._add_background(slide)

        # Title
        textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(1.5))
        text_frame = textbox.text_frame
        text_frame.text = title
        text_frame.paragraphs[0].font.name = 'Arial'
        text_frame.paragraphs[0].font.size = Pt(48)
        text_frame.paragraphs[0].font.color.rgb = self.COLORS['accent']
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtitle and contact
        if subtitle or contact:
            content = []
            if subtitle:
                content.append(subtitle)
            if contact:
                content.append(f"Contact: {contact}")

            textbox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(2))
            text_frame2 = textbox2.text_frame
            text_frame2.word_wrap = True

            for item in content:
                paragraph = text_frame2.paragraphs[0] if len(text_frame2.paragraphs) == 0 else text_frame2.add_paragraph()
                paragraph.text = item
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(28)
                paragraph.font.color.rgb = self.COLORS['text_secondary']
                paragraph.alignment = PP_ALIGN.CENTER

    def add_image_slide(self, title: str, image_path: str, caption: str = None):
        """
        Add a slide with a full-width image and optional caption.

        Args:
            title: Slide title
            image_path: Path to the image file
            caption: Optional caption below the image
        """
        from pathlib import Path
        
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])  # Blank layout
        self._add_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.name = 'Arial'
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = self.COLORS['text_primary']
        title_frame.paragraphs[0].font.bold = True

        # Image
        img_path = Path(image_path)
        if img_path.exists():
            # Calculate image dimensions to fit slide (max 11" wide, 5" tall)
            slide.shapes.add_picture(
                str(img_path),
                Inches(1.17),  # Center horizontally
                Inches(1.3),
                width=Inches(11),
                height=Inches(5.2)
            )
        else:
            # Placeholder if image not found
            placeholder_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
            placeholder_frame = placeholder_box.text_frame
            placeholder_frame.text = f"[Image not found: {img_path.name}]"
            placeholder_frame.paragraphs[0].font.color.rgb = self.COLORS['text_secondary']
            placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Caption
        if caption:
            caption_box = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11), Inches(0.5))
            caption_frame = caption_box.text_frame
            caption_frame.text = caption
            caption_frame.paragraphs[0].font.name = 'Arial'
            caption_frame.paragraphs[0].font.size = Pt(14)
            caption_frame.paragraphs[0].font.color.rgb = self.COLORS['text_secondary']
            caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_key_metrics_slide(self, title: str, metrics: dict, columns: int = 3):
        """
        Add a slide displaying key-value metrics in a professional grid layout.

        Args:
            title: Slide title
            metrics: Dictionary of metric_name: value pairs
            columns: Number of columns in the grid (default: 3)
        """
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[6])  # Blank layout
        self._add_background(slide)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.name = 'Arial'
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.color.rgb = self.COLORS['text_primary']
        title_frame.paragraphs[0].font.bold = True

        # Calculate grid layout
        items = list(metrics.items())
        rows = (len(items) + columns - 1) // columns
        
        cell_width = 11.5 / columns
        cell_height = 5.5 / max(rows, 1)
        start_x = 0.9
        start_y = 1.3

        for i, (key, value) in enumerate(items):
            col = i % columns
            row = i // columns
            
            x = start_x + col * cell_width
            y = start_y + row * cell_height

            # Metric box
            metric_box = slide.shapes.add_textbox(
                Inches(x), Inches(y),
                Inches(cell_width - 0.2), Inches(cell_height - 0.2)
            )
            frame = metric_box.text_frame
            frame.word_wrap = True

            # Metric name
            p1 = frame.paragraphs[0]
            p1.text = str(key)
            p1.font.name = 'Arial'
            p1.font.size = Pt(14)
            p1.font.color.rgb = self.COLORS['text_secondary']
            p1.alignment = PP_ALIGN.CENTER

            # Metric value
            p2 = frame.add_paragraph()
            p2.text = str(value)
            p2.font.name = 'Arial'
            p2.font.size = Pt(24)
            p2.font.color.rgb = self.COLORS['accent']
            p2.font.bold = True
            p2.alignment = PP_ALIGN.CENTER

    def save(self, filename: str = None):
        """
        Save the presentation.

        Args:
            filename: Optional custom filename
        """
        output_filename = filename or self.filename
        self.presentation.save(output_filename)
        print(f"Presentation saved as: {output_filename}")


# Example usage for your research presentation
if __name__ == "__main__":
    ppt = ResearchPresentationTemplate('optical_sensors.pptx')

    ppt.add_title_slide(
        title='Optical Fiber Sensors',
        subtitle='AI-Based Signal Processing for Real-Time Monitoring',
        author='Your Name',
        date='December 2025'
    )

    ppt.add_section_slide('Introduction')

    ppt.add_content_slide(
        title='Research Focus',
        content_items=[
            'FBG sensors for strain and temperature monitoring',
            'SPR-based optical sensors for biomedical applications',
            'Real-time ML-based demodulation',
            'Industrial IoT integration'
        ]
    )

    ppt.add_section_slide('Methodology')

    ppt.add_two_column_slide(
        title='Experimental vs. Simulation',
        left_content=[
            'Experimental Setup:',
            '• Broadband light source',
            '• Optical spectrum analyzer',
            '• 100 Hz sampling rate'
        ],
        right_content=[
            'Simulation:',
            '• COMSOL-based modeling',
            '• FEM analysis',
            '• Neural network training'
        ]
    )

    ppt.add_closing_slide(
        title='Thank You!',
        subtitle='Questions?',
        contact='your.email@chulalongkorn.ac.th'
    )

    ppt.save()
